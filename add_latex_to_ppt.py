import sys
import shutil
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def render_latex(formula, filename, fontsize=20):
    fig = plt.figure(figsize=(0.01, 0.01))
    fig.text(0, 0, f"${formula}$", fontsize=fontsize)
    fig.savefig(filename, dpi=300, transparent=True, format='png', bbox_inches='tight', pad_inches=0.1)
    plt.close(fig)

def remove_old_textboxes(slide, identifying_text):
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame and identifying_text in shape.text:
            shapes_to_remove.append(shape)
            
    for shape in shapes_to_remove:
        element = shape.element
        element.getparent().remove(element)

def add_math1_to_slide(slide):
    # Render formulas
    render_latex(r"EAR = \frac{||p_2 - p_6|| + ||p_3 - p_5||}{2 ||p_1 - p_4||}", "form1.png")
    render_latex(r"s \mathbf{p}_c = \mathbf{K} [\mathbf{R} | \mathbf{t}] \mathbf{P}_w", "form2.png")
    
    # Text 1
    t1 = "1. Eye Aspect Ratio (EAR) for Blink Detection:\nUsed to calculate Mb (Blink Modifier) for fatigue tracking.\n• Formula:"
    tx1 = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    p1 = tx1.text_frame.paragraphs[0]
    p1.text = t1
    p1.font.size = Pt(16)
    
    # Image 1
    slide.shapes.add_picture("form1.png", Inches(1.5), Inches(2.3), height=Inches(0.6))
    
    # Text 1 continued
    t1_cont = "• Where: p1, p4 are horizontal eye landmarks; p2, p3, p5, p6 are vertical eye landmarks."
    tx1_cont = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(8), Inches(0.5))
    p1_cont = tx1_cont.text_frame.paragraphs[0]
    p1_cont.text = t1_cont
    p1_cont.font.size = Pt(16)
    
    # Text 2
    t2 = "2. Head Pose Estimation (Perspective-n-Point / PnP):\nUsed to map 3D face model points to 2D image coordinates to find Pitch (θ), Yaw (φ), and Roll (ψ).\n• Formula:"
    tx2 = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(1))
    p2 = tx2.text_frame.paragraphs[0]
    p2.text = t2
    p2.font.size = Pt(16)
    
    # Image 2
    slide.shapes.add_picture("form2.png", Inches(1.5), Inches(4.5), height=Inches(0.4))
    
    # Text 2 continued
    t2_cont = "• Where:\n   s = scale factor\n   p_c = 2D image coordinates\n   K = Camera Intrinsic Matrix\n   [R | t] = Rotation and Translation matrix\n   P_w = 3D World coordinates"
    tx2_cont = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(8), Inches(1.5))
    p2_cont = tx2_cont.text_frame.paragraphs[0]
    p2_cont.text = t2_cont
    p2_cont.font.size = Pt(16)

def add_math2_to_slide(slide):
    # Render formula
    render_latex(r"E = \max(0, \min(100, B_e - P_h + M_b))", "form3.png")
    
    # Text 3
    t3 = "3. Total Engagement / Focus Score (E):\nCalculated on a scale of 0 to 100.\n• Formula:"
    tx3 = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    p3 = tx3.text_frame.paragraphs[0]
    p3.text = t3
    p3.font.size = Pt(16)
    
    # Image 3
    slide.shapes.add_picture("form3.png", Inches(1.5), Inches(2.3), height=Inches(0.4))
    
    # Text 3 continued
    t3_cont = "• Where:\n   B_e (Base Emotion Score): 85-95 for Happy/Surprise, 65-80 for Neutral, 30-50 for Negative.\n   P_h (Head Pose Penalty): -30 IF |φ| > 30° OR |θ| > 30° (Looking away).\n   M_b (Blink Modifier): -15 if blink rate > 35 (Fatigued), +10 if blink rate < 5 (Intense focus)."
    tx3_cont = slide.shapes.add_textbox(Inches(1), Inches(2.9), Inches(8), Inches(2))
    p3_cont = tx3_cont.text_frame.paragraphs[0]
    p3_cont.text = t3_cont
    p3_cont.font.size = Pt(16)

def main(input_path, output_path):
    shutil.copyfile(input_path, output_path)
    prs = Presentation(output_path)
    
    slide10 = prs.slides[9]
    slide11 = prs.slides[10]
    
    # Remove the plain text versions added in V2
    remove_old_textboxes(slide10, "1. Eye Aspect Ratio")
    remove_old_textboxes(slide11, "3. Total Engagement")
    
    # Add new text and latex images
    add_math1_to_slide(slide10)
    add_math2_to_slide(slide11)
                
    prs.save(output_path)
    print(f"Updated presentation saved to {output_path}")

if __name__ == '__main__':
    main(sys.argv[1], sys.argv[2])
