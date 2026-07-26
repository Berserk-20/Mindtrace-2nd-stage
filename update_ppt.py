import sys
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def replace_text_in_runs(shape_or_cell, replacements):
    if not hasattr(shape_or_cell, "text_frame") or shape_or_cell.text_frame is None:
        return
    for paragraph in shape_or_cell.text_frame.paragraphs:
        for run in paragraph.runs:
            for old, new in replacements.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    
        still_needs_replace = False
        for old in replacements.keys():
            if old in paragraph.text:
                still_needs_replace = True
                break
                
        if still_needs_replace:
            full_text = paragraph.text
            for old, new in replacements.items():
                full_text = full_text.replace(old, new)
            
            if len(paragraph.runs) > 0:
                paragraph.runs[0].text = full_text
                for r in paragraph.runs[1:]:
                    r.text = ""

def update_math_models(slide, content):
    # Find and remove the group shape or any shape that isn't title or footer
    # Usually shape 2 is GROUP (6) as seen in inspection
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.shape_type == 6: # GROUP
            shapes_to_remove.append(shape)
            
    for shape in shapes_to_remove:
        element = shape.element
        element.getparent().remove(element)
        
    # Add new text box with the math content
    # Coordinates approx: left=1, top=1.5, width=8, height=5
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.LEFT

def main(input_path, output_path):
    shutil.copyfile(input_path, output_path)
    prs = Presentation(output_path)
    
    replacements = {
        "Mini-Xception model": "ResNet18 model",
        "Mini-Xception CNN": "ResNet18 CNN",
        "Mini-Xception": "ResNet18",
        "MiniXception": "ResNet18",
        "Streamlit dashboard": "React dashboard",
        "Streamlit": "React",
        "48×48": "96x96",
        "48, 48, 1": "96, 96, 3",
        "TensorFlow, Keras, OpenCV, MediaPipe, Streamlit, NumPy, Pandas": "PyTorch, OpenCV, MediaPipe, FastAPI, React, Node.js, MongoDB",
    }
    
    # Slides to update (0-indexed): 9, 10, 11, 12, 13, 16
    slides_to_update = [9, 10, 11, 12, 13, 16]
    
    for i in slides_to_update:
        slide = prs.slides[i]
        
        # Specifically update text in these slides
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_text_in_runs(cell, replacements)
            else:
                replace_text_in_runs(shape, replacements)
                
    # Update Slide 10 (Math Model 1)
    math1 = (
        "1. Eye Aspect Ratio (EAR) for Blink Detection:\n"
        "Used to calculate Mb (Blink Modifier) for fatigue tracking.\n"
        "• Formula: EAR = (||p_2 - p_6|| + ||p_3 - p_5||) / (2 ||p_1 - p_4||)\n"
        "• Where: p_1, p_4 are horizontal eye landmarks; p_2, p_3, p_5, p_6 are vertical eye landmarks.\n\n"
        "2. Head Pose Estimation (Perspective-n-Point / PnP):\n"
        "Used to map 3D face model points to 2D image coordinates to find Pitch (θ), Yaw (φ), and Roll (ψ).\n"
        "• Formula: s p_c = K [R | t] P_w\n"
        "• Where:\n"
        "   s = scale factor\n"
        "   p_c = 2D image coordinates\n"
        "   K = Camera Intrinsic Matrix\n"
        "   [R | t] = Rotation and Translation matrix\n"
        "   P_w = 3D World coordinates"
    )
    update_math_models(prs.slides[9], math1)
    
    # Update Slide 11 (Math Model 2)
    math2 = (
        "3. Total Engagement / Focus Score (E):\n"
        "Calculated on a scale of 0 to 100.\n"
        "• Formula: E = max(0, min(100, B_e - P_h + M_b))\n"
        "• Where:\n"
        "   B_e (Base Emotion Score): 85-95 for Happy/Surprise, 65-80 for Neutral, 30-50 for Negative.\n"
        "   P_h (Head Pose Penalty): -30 IF |φ| > 30° OR |θ| > 30° (Looking away).\n"
        "   M_b (Blink Modifier): -15 if blink rate > 35 (Fatigued), +10 if blink rate < 5 (Intense focus)."
    )
    update_math_models(prs.slides[10], math2)
                
    prs.save(output_path)
    print(f"Updated presentation saved to {output_path}")

if __name__ == '__main__':
    main(sys.argv[1], sys.argv[2])
