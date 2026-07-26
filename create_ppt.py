import re
from pptx import Presentation
from pptx.util import Inches, Pt

def markdown_to_ppt(md_filepath, output_filepath):
    # Initialize presentation
    prs = Presentation()
    
    # Read markdown
    with open(md_filepath, 'r', encoding='utf-8') as f:
        content = f.read()
        
    # Split by slides
    # Each slide starts with '### Slide X: Title'
    slides_data = re.split(r'### Slide \d+: ', content)
    
    for slide_data in slides_data:
        slide_data = slide_data.strip()
        if not slide_data:
            continue
            
        # The first line is the title, the rest is content
        lines = slide_data.split('\n')
        title_text = lines[0].strip()
        body_text_lines = [line.strip() for line in lines[1:] if line.strip() and not line.startswith('---')]
        
        # If it's the "Introduction" or something before slide 1, we handle it if needed
        # But our regex splits specifically on '### Slide \d+: '
        
        # Create a slide
        # Layout 1 is title and content
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = title_text
        
        # Set body
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        
        first_paragraph = True
        
        for line in body_text_lines:
            # Clean up bullet points
            cleaned_line = line.strip()
            level = 0
            
            # Check for sub-bullets
            if cleaned_line.startswith('* '):
                cleaned_line = cleaned_line[2:]
            elif cleaned_line.startswith('- '):
                cleaned_line = cleaned_line[2:]
            elif line.startswith('  * ') or line.startswith('  - '):
                level = 1
                cleaned_line = cleaned_line.strip()[2:]
            elif re.match(r'^\d+\.', cleaned_line):
                # Clean numbered lists
                cleaned_line = re.sub(r'^\d+\.\s*', '', cleaned_line)
                
            # Clean markdown bold/italics
            cleaned_line = cleaned_line.replace('**', '').replace('*', '').replace('_', '')
            
            if first_paragraph:
                p = tf.paragraphs[0]
                first_paragraph = False
            else:
                p = tf.add_paragraph()
                
            p.text = cleaned_line
            p.level = level
            p.font.size = Pt(18)
            
    prs.save(output_filepath)
    print(f"Presentation saved to {output_filepath}")

if __name__ == "__main__":
    md_file = r"c:\Users\sanka\MindTrace\MindTrace_PPT_Content.md"
    out_file = "MindTrace_Presentation.pptx"
    markdown_to_ppt(md_file, out_file)
